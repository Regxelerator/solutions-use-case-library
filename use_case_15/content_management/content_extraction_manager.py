import os
import json
import io
import shutil
import base64
from typing import Any, Optional, Tuple, Dict

from docx import Document
from odf.opendocument import load
from odf.text import P
from odf import teletype
import pandas as pd
import pytesseract
import fitz                              
from PIL import Image
from pptx import Presentation

from utils.utils import (
    extract_images_from_pdf_content,
    analyze_image_with_openai,
    is_scanned_pdf,
    tesseract_exists,
)


class ContentManager:
    """
    Download a SharePoint file, extract its textual content (native or OCR),
    and (optionally) write derived artefacts back to SharePoint.
    """

    def __init__(self, graph_client: Any, output_drive_id: str) -> None:
        self.graph           = graph_client
        self.headers         = graph_client.headers
        self.output_drive_id = output_drive_id


    def get_file_content(
        self,
        drive_id: str,
        file_id: str,
        mime_type: str,
    ) -> Tuple[Optional[bytes], Optional[str]]:
        """
        Download a file from SharePoint and return its raw bytes plus extracted
        text (if any).
        """
        try:
            url  = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{file_id}/content"
            resp = self.graph.get(url)
            resp.raise_for_status()

            content        = resp.content
            extracted_text = self._extract_text_from_content(content, mime_type)
            return content, extracted_text
        except Exception as exc:
            print(f"Error fetching file {file_id}: {exc}")
            return None, None

    def get_text_content(
        self,
        file_obj: Dict[str, Any],
    ) -> Tuple[Optional[bytes], Optional[str]]:
        """
        Convenience wrapper that extracts the Graph IDs we need from an item
        JSON blob and calls `get_file_content`.
        """
        try:
            drive_id  = file_obj.get("parentReference", {}).get("driveId")
            file_id   = file_obj.get("id")
            mime_type = file_obj.get("file", {}).get("mimeType")
            if not all((drive_id, file_id, mime_type)):
                return None, None
            return self.get_file_content(drive_id, file_id, mime_type)
        except Exception as exc:
            print(f"Error resolving drive/file IDs: {exc}")
            return None, None


    def _extract_text_from_content(
        self,
        content: bytes,
        mime_type: str,
    ) -> Optional[str]:
        """
        Route to the correct extractor based on MIME type.
        """
        try:
            if mime_type == "application/pdf":
                return self._extract_pdf_text(content)

            if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return self._extract_docx_text(content)

            if mime_type in {
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
                "application/vnd.ms-excel.sheet.macroEnabled.12",
            }:
                return self._extract_excel_text(content)

            if mime_type == "application/vnd.oasis.opendocument.text":
                return self._extract_odt_text(content)

            if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return self._extract_text_from_pptx(content)

            if mime_type.startswith("image/"):
                image = Image.open(io.BytesIO(content))
                buf   = io.BytesIO()
                image.save(buf, format="PNG")
                encoded = base64.b64encode(buf.getvalue()).decode()
                md      = analyze_image_with_openai(encoded)
                return md.get("image_description") if isinstance(md, dict) else None

            print(f"Unsupported MIME type for text extraction: {mime_type}")
            return None
        except Exception as exc:
            print(f"Text-extraction error: {exc}")
            return None


    @staticmethod
    def _extract_odt_text(content: bytes) -> Optional[str]:
        try:
            doc   = load(io.BytesIO(content))
            paras = doc.getElementsByType(P)
            return "\n".join(teletype.extractText(p) for p in paras).strip()
        except Exception as exc:
            print(f"ODT extraction error: {exc}")
            return None

    def _extract_pdf_text(self, content: bytes) -> str:
        """
        Combine embedded text and OCR for scanned pages (if Tesseract present).
        """
        full_text = []
        doc       = fitz.open(stream=content, filetype="pdf")
        scannable, _ratio = is_scanned_pdf(content)

        for page in doc:
            txt = page.get_text("text").strip()
            if txt:
                full_text.append(txt)
                continue

            if not (scannable and tesseract_exists()):
                continue

            if page.get_images() or page.get_drawings():
                pix = page.get_pixmap(dpi=300)
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                try:
                    ocr = pytesseract.image_to_string(img).strip()
                    if ocr:
                        full_text.append(ocr)
                except pytesseract.TesseractNotFoundError:
                    pass

        return "\n".join(full_text)

    @staticmethod
    def _extract_docx_text(content: bytes) -> str:
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs)

    @staticmethod
    def _extract_excel_text(content: bytes) -> str:
        df = pd.read_excel(io.BytesIO(content))
        return df.to_string(index=False)

    @staticmethod
    def _extract_text_from_pptx(content: bytes) -> Optional[str]:
        try:
            prs = Presentation(io.BytesIO(content))
            return "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )
        except Exception as exc:
            print(f"PPTX extraction error: {exc}")
            return None


    def process_and_upload_file(
        self,
        output_drive_id: str,
        drive_id: str,
        source_file_id: str,
        target_folder_id: str,
        metadata: Dict[str, Any],
    ) -> Dict[str, Any]:
        """
        Download the source file, build a consolidated `content.json`
        (metadata + extracted text), optionally extract embedded images,
        and upload everything to the target folder.
        """
        try:
            url  = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{source_file_id}/content"
            resp = self.graph.get(url)
            resp.raise_for_status()
            file_content = resp.content

            mime_type   = metadata.get("mimeType")
            text        = metadata.get("content_text")
            combined: Dict[str, Any] = {}

            if mime_type and not mime_type.startswith("image/"):
                if text is None:
                    text = self._extract_text_from_content(file_content, mime_type)
                combined["file_name"] = metadata.get("name")
                combined.update(
                    {k: v for k, v in metadata.items() if k not in {"name", "content_text", "subfolderPath", "subfolder"}}
                )
                combined["text"] = text
            else:   # pure image
                combined["file_name"] = metadata.get("name")
                combined.update(
                    {k: v for k, v in metadata.items() if k not in {"name", "content_text", "subfolderPath", "subfolder"}}
                )
                combined["text"] = metadata.get("img_text")

            combined_bytes = json.dumps(combined, indent=2).encode("utf-8")
            content_file_id = self.graph.upload_file_to_folder(
                output_drive_id, target_folder_id, "content.json", combined_bytes
            )

            if (
                not metadata.get("is_pdf_scanned", False)
                and mime_type
                and mime_type.startswith("application/pdf")
            ):
                temp_dir, img_count, img_meta = extract_images_from_pdf_content(file_content)

                if img_meta:
                    self.graph.upload_file_to_folder(
                        output_drive_id,
                        target_folder_id,
                        "images.json",
                        json.dumps(img_meta, indent=2).encode("utf-8"),
                    )

                try:
                    if img_count:
                        sub_id = self.graph.create_subfolder(output_drive_id, target_folder_id, "images")
                        for img_name in os.listdir(temp_dir):
                            with open(os.path.join(temp_dir, img_name), "rb") as fh:
                                self.graph.upload_file_to_folder(output_drive_id, sub_id, img_name, fh.read())
                finally:
                    shutil.rmtree(temp_dir)

            return {"content_file_id": content_file_id}
        except Exception as exc:
            print(f"process_and_upload_file error: {exc}")
            raise