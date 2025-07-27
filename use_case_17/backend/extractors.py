from pathlib import Path

def _from_pdf(path: Path) -> str:
    import PyPDF2 
    text = ""
    with open(path, "rb") as fh:
        reader = PyPDF2.PdfReader(fh)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text


def _from_docx(path: Path) -> str:
    import docx  
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _from_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)
    return "\n".join(lines)


def extract_text_from_file(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _from_pdf(path)
    if suffix == ".docx":
        return _from_docx(path)
    if suffix == ".pptx":
        return _from_pptx(path)
    if suffix in {".txt", ".text"}:
        return path.read_text(encoding="utf-8", errors="ignore")

    raise RuntimeError(f"Unsupported file type: {suffix}")
